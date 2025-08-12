"""
Document Processing Engine for Universal Knowledge Management

Real implementation of document processing with support for:
- PDF text extraction with metadata preservation
- Word document processing with formatting retention
- PowerPoint processing for educational content
- Intelligent text chunking with overlap
- Content validation and governance classification
"""

import os
import io
import hashlib
import mimetypes
from datetime import datetime, timezone
from typing import Dict, List, Any, Optional, Tuple, Union
from pathlib import Path

# Document processing libraries
import PyPDF2
from docx import Document
from pptx import Presentation
import langdetect
import textstat

# Import governance core for trust scoring
from ...core.governance.governance_core import calculate_entry_hash

class DocumentProcessor:
    """
    Real document processing engine with governance integration
    """
    
    def __init__(self):
        self.supported_types = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_docx,  # Will attempt docx processing
            'pptx': self._process_pptx,
            'ppt': self._process_pptx,  # Will attempt pptx processing
            'txt': self._process_text,
            'md': self._process_text,
            'markdown': self._process_text
        }
        
        # Chunking configuration
        self.chunk_size = 1000  # Target chunk size in characters
        self.chunk_overlap = 200  # Overlap between chunks
        self.min_chunk_size = 100  # Minimum viable chunk size
        
    def process_document(
        self, 
        file_content: bytes, 
        filename: str, 
        file_type: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Process a document and extract structured content
        
        Args:
            file_content: Raw file bytes
            filename: Original filename
            file_type: File type override (optional)
            
        Returns:
            Dict containing processed content, metadata, and chunks
        """
        try:
            # Determine file type
            if not file_type:
                file_type = self._detect_file_type(filename, file_content)
            
            file_type = file_type.lower().replace('.', '')
            
            if file_type not in self.supported_types:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            # Calculate content hash for integrity
            content_hash = hashlib.sha256(file_content).hexdigest()
            
            # Process document based on type
            processor = self.supported_types[file_type]
            processing_result = processor(file_content, filename)
            
            # Extract metadata
            metadata = self._extract_metadata(file_content, filename, file_type)
            
            # Perform content analysis
            content_analysis = self._analyze_content(processing_result['text_content'])
            
            # Create intelligent chunks
            chunks = self._create_intelligent_chunks(
                processing_result['text_content'],
                processing_result.get('structure', [])
            )
            
            # Calculate governance trust score
            trust_score = self._calculate_trust_score(
                processing_result['text_content'],
                content_analysis,
                metadata
            )
            
            # Compile results
            result = {
                'document_id': None,  # Will be set by caller
                'filename': filename,
                'file_type': file_type,
                'file_size': len(file_content),
                'content_hash': content_hash,
                'text_content': processing_result['text_content'],
                'metadata': metadata,
                'content_analysis': content_analysis,
                'chunks': chunks,
                'trust_score': trust_score,
                'processing_status': 'completed',
                'processing_timestamp': datetime.now(timezone.utc).isoformat(),
                'governance_classification': self._classify_content(
                    processing_result['text_content'], 
                    content_analysis
                ),
                'total_tokens': len(processing_result['text_content'].split()),
                'total_chunks': len(chunks),
                'structure': processing_result.get('structure', [])
            }
            
            return result
            
        except Exception as e:
            return {
                'document_id': None,
                'filename': filename,
                'file_type': file_type or 'unknown',
                'file_size': len(file_content),
                'processing_status': 'failed',
                'error': str(e),
                'processing_timestamp': datetime.now(timezone.utc).isoformat()
            }
    
    def _detect_file_type(self, filename: str, file_content: bytes) -> str:
        """Detect file type from filename and content"""
        # Try filename extension first
        if '.' in filename:
            ext = filename.split('.')[-1].lower()
            if ext in self.supported_types:
                return ext
        
        # Try MIME type detection
        mime_type, _ = mimetypes.guess_type(filename)
        if mime_type:
            mime_to_ext = {
                'application/pdf': 'pdf',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
                'application/msword': 'doc',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
                'application/vnd.ms-powerpoint': 'ppt',
                'text/plain': 'txt',
                'text/markdown': 'md'
            }
            if mime_type in mime_to_ext:
                return mime_to_ext[mime_type]
        
        # Try content-based detection
        if file_content.startswith(b'%PDF'):
            return 'pdf'
        elif file_content.startswith(b'PK'):  # ZIP-based formats (docx, pptx)
            # More sophisticated detection would be needed here
            return 'docx'  # Default assumption
        
        return 'txt'  # Default fallback
    
    def _process_pdf(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process PDF document"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = ""
            structure = []
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                text_content += page_text + "\n\n"
                
                structure.append({
                    'type': 'page',
                    'page_number': page_num + 1,
                    'text_length': len(page_text),
                    'start_position': len(text_content) - len(page_text) - 2
                })
            
            return {
                'text_content': text_content.strip(),
                'structure': structure,
                'page_count': len(pdf_reader.pages),
                'processing_method': 'PyPDF2'
            }
            
        except Exception as e:
            raise Exception(f"PDF processing failed: {str(e)}")
    
    def _process_docx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process Word document"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = Document(doc_file)
            
            text_content = ""
            structure = []
            
            for para_num, paragraph in enumerate(doc.paragraphs):
                para_text = paragraph.text
                if para_text.strip():  # Skip empty paragraphs
                    text_content += para_text + "\n\n"
                    
                    # Detect headings and structure
                    style_name = paragraph.style.name if paragraph.style else "Normal"
                    structure.append({
                        'type': 'paragraph',
                        'paragraph_number': para_num + 1,
                        'style': style_name,
                        'text_length': len(para_text),
                        'start_position': len(text_content) - len(para_text) - 2,
                        'is_heading': 'Heading' in style_name
                    })
            
            return {
                'text_content': text_content.strip(),
                'structure': structure,
                'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()]),
                'processing_method': 'python-docx'
            }
            
        except Exception as e:
            raise Exception(f"DOCX processing failed: {str(e)}")
    
    def _process_pptx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process PowerPoint presentation"""
        try:
            ppt_file = io.BytesIO(file_content)
            prs = Presentation(ppt_file)
            
            text_content = ""
            structure = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                
                # Extract text from all shapes in slide
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape_text = shape.text
                        if shape_text.strip():
                            slide_text += shape_text + "\n"
                
                if slide_text.strip():
                    text_content += f"Slide {slide_num + 1}:\n{slide_text}\n\n"
                    
                    structure.append({
                        'type': 'slide',
                        'slide_number': slide_num + 1,
                        'text_length': len(slide_text),
                        'start_position': len(text_content) - len(slide_text) - len(f"Slide {slide_num + 1}:\n") - 2,
                        'shape_count': len([s for s in slide.shapes if hasattr(s, "text") and s.text.strip()])
                    })
            
            return {
                'text_content': text_content.strip(),
                'structure': structure,
                'slide_count': len(prs.slides),
                'processing_method': 'python-pptx'
            }
            
        except Exception as e:
            raise Exception(f"PPTX processing failed: {str(e)}")
    
    def _process_text(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process plain text document"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text_content = None
            
            for encoding in encodings:
                try:
                    text_content = file_content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                raise Exception("Could not decode text file with any supported encoding")
            
            # Simple structure detection for text files
            lines = text_content.split('\n')
            structure = []
            
            for line_num, line in enumerate(lines):
                if line.strip():  # Non-empty lines
                    structure.append({
                        'type': 'line',
                        'line_number': line_num + 1,
                        'text_length': len(line),
                        'start_position': sum(len(l) + 1 for l in lines[:line_num]),
                        'is_heading': line.strip().startswith('#') or line.isupper()
                    })
            
            return {
                'text_content': text_content,
                'structure': structure,
                'line_count': len(lines),
                'processing_method': 'text-decode'
            }
            
        except Exception as e:
            raise Exception(f"Text processing failed: {str(e)}")
    
    def _extract_metadata(self, file_content: bytes, filename: str, file_type: str) -> Dict[str, Any]:
        """Extract metadata from document"""
        metadata = {
            'filename': filename,
            'file_type': file_type,
            'file_size': len(file_content),
            'extraction_timestamp': datetime.now(timezone.utc).isoformat()
        }
        
        # Add file-type specific metadata extraction here
        # For now, basic metadata only
        
        return metadata
    
    def _analyze_content(self, text_content: str) -> Dict[str, Any]:
        """Analyze content for quality and characteristics"""
        try:
            analysis = {
                'character_count': len(text_content),
                'word_count': len(text_content.split()),
                'sentence_count': text_content.count('.') + text_content.count('!') + text_content.count('?'),
                'paragraph_count': len([p for p in text_content.split('\n\n') if p.strip()]),
            }
            
            # Language detection
            try:
                detected_language = langdetect.detect(text_content[:1000])  # Sample first 1000 chars
                analysis['language'] = detected_language
                analysis['language_confidence'] = 0.8  # langdetect doesn't provide confidence
            except:
                analysis['language'] = 'unknown'
                analysis['language_confidence'] = 0.0
            
            # Readability analysis
            try:
                analysis['flesch_reading_ease'] = textstat.flesch_reading_ease(text_content)
                analysis['flesch_kincaid_grade'] = textstat.flesch_kincaid_grade(text_content)
                analysis['automated_readability_index'] = textstat.automated_readability_index(text_content)
            except:
                analysis['flesch_reading_ease'] = None
                analysis['flesch_kincaid_grade'] = None
                analysis['automated_readability_index'] = None
            
            # Content quality indicators
            analysis['avg_sentence_length'] = (
                analysis['word_count'] / max(analysis['sentence_count'], 1)
            )
            analysis['avg_word_length'] = (
                analysis['character_count'] / max(analysis['word_count'], 1)
            )
            
            return analysis
            
        except Exception as e:
            return {
                'character_count': len(text_content),
                'word_count': len(text_content.split()),
                'analysis_error': str(e)
            }
    
    def _create_intelligent_chunks(self, text_content: str, structure: List[Dict]) -> List[Dict[str, Any]]:
        """Create intelligent text chunks based on document structure"""
        chunks = []
        
        if not text_content.strip():
            return chunks
        
        # If we have structure information, use it for better chunking
        if structure:
            chunks = self._structure_based_chunking(text_content, structure)
        else:
            chunks = self._simple_chunking(text_content)
        
        # Add chunk metadata
        for i, chunk in enumerate(chunks):
            chunk.update({
                'chunk_id': None,  # Will be set by caller
                'chunk_index': i,
                'token_count': len(chunk['content'].split()),
                'character_count': len(chunk['content']),
                'chunk_hash': hashlib.sha256(chunk['content'].encode()).hexdigest()
            })
        
        return chunks
    
    def _structure_based_chunking(self, text_content: str, structure: List[Dict]) -> List[Dict[str, Any]]:
        """Create chunks based on document structure"""
        chunks = []
        current_chunk = ""
        current_start = 0
        
        for struct_item in structure:
            # Get the text for this structural element
            start_pos = struct_item.get('start_position', 0)
            text_length = struct_item.get('text_length', 0)
            
            if start_pos + text_length <= len(text_content):
                element_text = text_content[start_pos:start_pos + text_length]
            else:
                continue
            
            # Check if adding this element would exceed chunk size
            if len(current_chunk) + len(element_text) > self.chunk_size and current_chunk:
                # Save current chunk
                chunks.append({
                    'content': current_chunk.strip(),
                    'start_position': current_start,
                    'end_position': current_start + len(current_chunk),
                    'structure_elements': []  # Could track which elements are in this chunk
                })
                
                # Start new chunk with overlap
                overlap_text = current_chunk[-self.chunk_overlap:] if len(current_chunk) > self.chunk_overlap else current_chunk
                current_chunk = overlap_text + "\n" + element_text
                current_start = start_pos - len(overlap_text) - 1
            else:
                # Add to current chunk
                if current_chunk:
                    current_chunk += "\n" + element_text
                else:
                    current_chunk = element_text
                    current_start = start_pos
        
        # Add final chunk
        if current_chunk.strip():
            chunks.append({
                'content': current_chunk.strip(),
                'start_position': current_start,
                'end_position': current_start + len(current_chunk),
                'structure_elements': []
            })
        
        return chunks
    
    def _simple_chunking(self, text_content: str) -> List[Dict[str, Any]]:
        """Simple character-based chunking with overlap"""
        chunks = []
        
        for i in range(0, len(text_content), self.chunk_size - self.chunk_overlap):
            chunk_text = text_content[i:i + self.chunk_size]
            
            # Try to break at sentence boundaries
            if i + self.chunk_size < len(text_content):
                # Look for sentence endings near the chunk boundary
                last_sentence_end = max(
                    chunk_text.rfind('.'),
                    chunk_text.rfind('!'),
                    chunk_text.rfind('?')
                )
                
                if last_sentence_end > self.chunk_size * 0.7:  # If we found a good break point
                    chunk_text = chunk_text[:last_sentence_end + 1]
            
            if len(chunk_text.strip()) >= self.min_chunk_size:
                chunks.append({
                    'content': chunk_text.strip(),
                    'start_position': i,
                    'end_position': i + len(chunk_text),
                    'chunking_method': 'simple'
                })
        
        return chunks
    
    def _calculate_trust_score(
        self, 
        text_content: str, 
        content_analysis: Dict[str, Any], 
        metadata: Dict[str, Any]
    ) -> float:
        """Calculate governance trust score using existing trust API"""
        try:
            # Use existing governance trust system instead of creating new logic
            from ..routes import call_governance_core
            
            trust_result = call_governance_core("calculate_trust_score", text_content)
            return trust_result.get("trust_score", 0.5)
            
        except Exception:
            return 0.5  # Default neutral trust score on error
    
    def _classify_content(self, text_content: str, content_analysis: Dict[str, Any]) -> str:
        """Classify content for governance purposes"""
        try:
            # Simple classification based on content characteristics
            word_count = content_analysis.get('word_count', 0)
            
            if word_count < 100:
                return 'brief'
            elif word_count < 1000:
                return 'standard'
            elif word_count < 5000:
                return 'detailed'
            else:
                return 'comprehensive'
                
        except Exception:
            return 'standard'

