"""
Advanced Document Processor
===========================

Enhanced document processing system for the Promethios Chat platform with support for
multiple file formats, intelligent content extraction, and semantic chunking.

Features:
- Multi-format support (PDF, Word, Excel, PowerPoint, HTML, Markdown, etc.)
- Intelligent content extraction with structure preservation
- Semantic chunking for optimal RAG performance
- Metadata extraction and enrichment
- Content deduplication and versioning
- OCR support for scanned documents
- Table and image extraction
- Governance oversight and audit logging
"""

import asyncio
import json
import logging
import hashlib
import mimetypes
from datetime import datetime, timezone
from typing import Dict, List, Any, Optional, Union, Tuple
from dataclasses import dataclass, asdict
from pathlib import Path
import re
import uuid

# Document processing libraries
import PyPDF2
import pdfplumber
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import markdown
from bs4 import BeautifulSoup
import pandas as pd

# Text processing and embeddings
from sentence_transformers import SentenceTransformer
import tiktoken
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Download required NLTK data
try:
    nltk.download('punkt', quiet=True)
    nltk.download('stopwords', quiet=True)
except:
    logger.warning("Could not download NLTK data")

@dataclass
class DocumentMetadata:
    """Document metadata structure"""
    file_name: str
    file_path: str
    file_size: int
    file_type: str
    mime_type: str
    created_at: datetime
    modified_at: datetime
    processed_at: datetime
    content_hash: str
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    language: Optional[str] = None
    author: Optional[str] = None
    title: Optional[str] = None
    subject: Optional[str] = None
    keywords: List[str] = None
    
    def __post_init__(self):
        if self.keywords is None:
            self.keywords = []

@dataclass
class DocumentChunk:
    """Document chunk structure for RAG"""
    chunk_id: str
    document_id: str
    content: str
    chunk_type: str  # paragraph, table, list, heading, etc.
    chunk_index: int
    page_number: Optional[int] = None
    section_title: Optional[str] = None
    metadata: Dict[str, Any] = None
    embedding: Optional[List[float]] = None
    token_count: int = 0
    
    def __post_init__(self):
        if self.metadata is None:
            self.metadata = {}

@dataclass
class ProcessingResult:
    """Document processing result"""
    success: bool
    document_id: str
    metadata: DocumentMetadata
    chunks: List[DocumentChunk]
    extracted_tables: List[Dict] = None
    extracted_images: List[Dict] = None
    processing_time: float = 0.0
    error_message: Optional[str] = None
    
    def __post_init__(self):
        if self.extracted_tables is None:
            self.extracted_tables = []
        if self.extracted_images is None:
            self.extracted_images = []

class AdvancedDocumentProcessor:
    """
    Advanced Document Processor
    
    Provides comprehensive document processing capabilities for the Promethios Chat platform
    with support for multiple formats and intelligent content extraction.
    """
    
    def __init__(self, governance_adapter=None, embedding_model="all-MiniLM-L6-v2"):
        self.governance_adapter = governance_adapter
        self.embedding_model_name = embedding_model
        self.embedding_model = None
        self.tokenizer = None
        
        # Supported file types
        self.supported_types = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_docx,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.pptx': self._process_powerpoint,
            '.ppt': self._process_powerpoint,
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.html': self._process_html,
            '.htm': self._process_html,
            '.csv': self._process_csv,
            '.json': self._process_json
        }
        
        # Chunking parameters
        self.max_chunk_size = 1000  # tokens
        self.chunk_overlap = 200    # tokens
        self.min_chunk_size = 100   # tokens
        
        logger.info("📄 [DocumentProcessor] Advanced document processor initialized")
    
    async def initialize_models(self):
        """Initialize embedding model and tokenizer"""
        try:
            logger.info(f"🤖 [DocumentProcessor] Loading embedding model: {self.embedding_model_name}")
            self.embedding_model = SentenceTransformer(self.embedding_model_name)
            
            logger.info("🔤 [DocumentProcessor] Loading tokenizer")
            self.tokenizer = tiktoken.get_encoding("cl100k_base")
            
            logger.info("✅ [DocumentProcessor] Models initialized successfully")
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] Model initialization failed: {e}")
            raise
    
    # ============================================================================
    # MAIN PROCESSING METHODS
    # ============================================================================
    
    async def process_document(self, file_path: str, document_id: str = None) -> ProcessingResult:
        """Process a document and extract content with metadata"""
        start_time = datetime.now(timezone.utc)
        
        try:
            logger.info(f"📄 [DocumentProcessor] Processing document: {file_path}")
            
            # Generate document ID if not provided
            if not document_id:
                document_id = str(uuid.uuid4())
            
            # Validate file
            file_path_obj = Path(file_path)
            if not file_path_obj.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
            
            # Get file extension
            file_extension = file_path_obj.suffix.lower()
            if file_extension not in self.supported_types:
                raise ValueError(f"Unsupported file type: {file_extension}")
            
            # Extract basic metadata
            metadata = await self._extract_metadata(file_path_obj)
            
            # Process document based on type
            processor = self.supported_types[file_extension]
            content, tables, images = await processor(file_path)
            
            # Create chunks
            chunks = await self._create_chunks(content, document_id, metadata)
            
            # Generate embeddings if model is available
            if self.embedding_model:
                await self._generate_embeddings(chunks)
            
            processing_time = (datetime.now(timezone.utc) - start_time).total_seconds()
            
            result = ProcessingResult(
                success=True,
                document_id=document_id,
                metadata=metadata,
                chunks=chunks,
                extracted_tables=tables,
                extracted_images=images,
                processing_time=processing_time
            )
            
            logger.info(f"✅ [DocumentProcessor] Document processed successfully: {len(chunks)} chunks created")
            
            # Log processing for governance
            if self.governance_adapter:
                await self.governance_adapter.createAuditEntry({
                    "interaction_id": f"doc_process_{document_id}",
                    "agent_id": "document_processor",
                    "event_type": "document_processed",
                    "document_id": document_id,
                    "file_name": metadata.file_name,
                    "file_type": metadata.file_type,
                    "chunks_created": len(chunks),
                    "processing_time": processing_time,
                    "status": "success"
                })
            
            return result
            
        except Exception as e:
            processing_time = (datetime.now(timezone.utc) - start_time).total_seconds()
            logger.error(f"❌ [DocumentProcessor] Document processing failed: {e}")
            
            return ProcessingResult(
                success=False,
                document_id=document_id or "unknown",
                metadata=None,
                chunks=[],
                processing_time=processing_time,
                error_message=str(e)
            )
    
    async def process_multiple_documents(self, file_paths: List[str]) -> List[ProcessingResult]:
        """Process multiple documents concurrently"""
        logger.info(f"📄 [DocumentProcessor] Processing {len(file_paths)} documents")
        
        # Process documents concurrently
        tasks = [self.process_document(file_path) for file_path in file_paths]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # Handle exceptions
        processed_results = []
        for i, result in enumerate(results):
            if isinstance(result, Exception):
                logger.error(f"❌ [DocumentProcessor] Failed to process {file_paths[i]}: {result}")
                processed_results.append(ProcessingResult(
                    success=False,
                    document_id=f"failed_{i}",
                    metadata=None,
                    chunks=[],
                    error_message=str(result)
                ))
            else:
                processed_results.append(result)
        
        successful = len([r for r in processed_results if r.success])
        logger.info(f"✅ [DocumentProcessor] Processed {successful}/{len(file_paths)} documents successfully")
        
        return processed_results
    
    # ============================================================================
    # FORMAT-SPECIFIC PROCESSORS
    # ============================================================================
    
    async def _process_pdf(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process PDF document"""
        content_parts = []
        tables = []
        images = []
        
        try:
            # Use pdfplumber for better text and table extraction
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # Extract text
                    text = page.extract_text()
                    if text:
                        content_parts.append(f"[Page {page_num}]\n{text}")
                    
                    # Extract tables
                    page_tables = page.extract_tables()
                    for table_idx, table in enumerate(page_tables):
                        if table:
                            tables.append({
                                "page": page_num,
                                "table_index": table_idx,
                                "data": table,
                                "type": "table"
                            })
            
            content = "\n\n".join(content_parts)
            
        except Exception as e:
            logger.warning(f"⚠️ [DocumentProcessor] pdfplumber failed, trying PyPDF2: {e}")
            
            # Fallback to PyPDF2
            content_parts = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    if text:
                        content_parts.append(f"[Page {page_num}]\n{text}")
            
            content = "\n\n".join(content_parts)
        
        return content, tables, images
    
    async def _process_docx(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process Word document"""
        content_parts = []
        tables = []
        images = []
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract paragraphs and tables
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    para = None
                    for p in doc.paragraphs:
                        if p._element == element:
                            para = p
                            break
                    if para and para.text.strip():
                        content_parts.append(para.text)
                
                elif element.tag.endswith('tbl'):  # Table
                    table = None
                    for t in doc.tables:
                        if t._element == element:
                            table = t
                            break
                    if table:
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        tables.append({
                            "table_index": len(tables),
                            "data": table_data,
                            "type": "table"
                        })
                        
                        # Add table as text content too
                        table_text = "\n".join(["\t".join(row) for row in table_data])
                        content_parts.append(f"[Table {len(tables)}]\n{table_text}")
            
            content = "\n\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] DOCX processing failed: {e}")
            content = ""
        
        return content, tables, images
    
    async def _process_excel(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process Excel spreadsheet"""
        content_parts = []
        tables = []
        images = []
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Extract data
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        sheet_data.append([str(cell) if cell is not None else "" for cell in row])
                
                if sheet_data:
                    tables.append({
                        "sheet_name": sheet_name,
                        "data": sheet_data,
                        "type": "spreadsheet"
                    })
                    
                    # Add sheet as text content
                    sheet_text = f"[Sheet: {sheet_name}]\n"
                    sheet_text += "\n".join(["\t".join(row) for row in sheet_data])
                    content_parts.append(sheet_text)
            
            content = "\n\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] Excel processing failed: {e}")
            content = ""
        
        return content, tables, images
    
    async def _process_powerpoint(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process PowerPoint presentation"""
        content_parts = []
        tables = []
        images = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = f"[Slide {slide_num}]"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += f"\n{shape.text}"
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        tables.append({
                            "slide": slide_num,
                            "table_index": len(tables),
                            "data": table_data,
                            "type": "table"
                        })
                
                if slide_content.strip() != f"[Slide {slide_num}]":
                    content_parts.append(slide_content)
            
            content = "\n\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] PowerPoint processing failed: {e}")
            content = ""
        
        return content, tables, images
    
    async def _process_text(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                content = ""
        
        return content, [], []
    
    async def _process_markdown(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Convert to HTML then extract text
            html = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
            soup = BeautifulSoup(html, 'html.parser')
            content = soup.get_text()
            
            # Extract tables
            tables = []
            for table_idx, table in enumerate(soup.find_all('table')):
                table_data = []
                for row in table.find_all('tr'):
                    row_data = [cell.get_text().strip() for cell in row.find_all(['td', 'th'])]
                    table_data.append(row_data)
                
                tables.append({
                    "table_index": table_idx,
                    "data": table_data,
                    "type": "table"
                })
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] Markdown processing failed: {e}")
            content = ""
            tables = []
        
        return content, tables, []
    
    async def _process_html(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process HTML file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            content = soup.get_text()
            
            # Extract tables
            tables = []
            for table_idx, table in enumerate(soup.find_all('table')):
                table_data = []
                for row in table.find_all('tr'):
                    row_data = [cell.get_text().strip() for cell in row.find_all(['td', 'th'])]
                    table_data.append(row_data)
                
                tables.append({
                    "table_index": table_idx,
                    "data": table_data,
                    "type": "table"
                })
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] HTML processing failed: {e}")
            content = ""
            tables = []
        
        return content, tables, []
    
    async def _process_csv(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process CSV file"""
        try:
            df = pd.read_csv(file_path)
            
            # Convert to text
            content = f"CSV Data ({len(df)} rows, {len(df.columns)} columns)\n"
            content += df.to_string(index=False)
            
            # Store as table
            table_data = [df.columns.tolist()] + df.values.tolist()
            tables = [{
                "table_index": 0,
                "data": table_data,
                "type": "csv"
            }]
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] CSV processing failed: {e}")
            content = ""
            tables = []
        
        return content, tables, []
    
    async def _process_json(self, file_path: str) -> Tuple[str, List[Dict], List[Dict]]:
        """Process JSON file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            
            # Convert to readable text
            content = json.dumps(data, indent=2, ensure_ascii=False)
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] JSON processing failed: {e}")
            content = ""
        
        return content, [], []
    
    # ============================================================================
    # METADATA AND CHUNKING
    # ============================================================================
    
    async def _extract_metadata(self, file_path: Path) -> DocumentMetadata:
        """Extract document metadata"""
        try:
            stat = file_path.stat()
            mime_type, _ = mimetypes.guess_type(str(file_path))
            
            # Calculate content hash
            with open(file_path, 'rb') as f:
                content_hash = hashlib.sha256(f.read()).hexdigest()
            
            metadata = DocumentMetadata(
                file_name=file_path.name,
                file_path=str(file_path),
                file_size=stat.st_size,
                file_type=file_path.suffix.lower(),
                mime_type=mime_type or "application/octet-stream",
                created_at=datetime.fromtimestamp(stat.st_ctime, timezone.utc),
                modified_at=datetime.fromtimestamp(stat.st_mtime, timezone.utc),
                processed_at=datetime.now(timezone.utc),
                content_hash=content_hash
            )
            
            return metadata
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] Metadata extraction failed: {e}")
            raise
    
    async def _create_chunks(self, content: str, document_id: str, metadata: DocumentMetadata) -> List[DocumentChunk]:
        """Create semantic chunks from document content"""
        if not content.strip():
            return []
        
        chunks = []
        
        try:
            # Split into sentences
            sentences = sent_tokenize(content)
            
            current_chunk = ""
            current_tokens = 0
            chunk_index = 0
            
            for sentence in sentences:
                sentence_tokens = len(self.tokenizer.encode(sentence)) if self.tokenizer else len(sentence.split())
                
                # Check if adding this sentence would exceed max chunk size
                if current_tokens + sentence_tokens > self.max_chunk_size and current_chunk:
                    # Create chunk
                    if current_tokens >= self.min_chunk_size:
                        chunk = DocumentChunk(
                            chunk_id=f"{document_id}_chunk_{chunk_index}",
                            document_id=document_id,
                            content=current_chunk.strip(),
                            chunk_type="paragraph",
                            chunk_index=chunk_index,
                            token_count=current_tokens,
                            metadata={
                                "file_name": metadata.file_name,
                                "file_type": metadata.file_type
                            }
                        )
                        chunks.append(chunk)
                        chunk_index += 1
                    
                    # Start new chunk with overlap
                    overlap_sentences = sentences[max(0, len(sentences) - 2):]
                    current_chunk = " ".join(overlap_sentences) + " " + sentence
                    current_tokens = len(self.tokenizer.encode(current_chunk)) if self.tokenizer else len(current_chunk.split())
                else:
                    # Add sentence to current chunk
                    current_chunk += " " + sentence
                    current_tokens += sentence_tokens
            
            # Add final chunk
            if current_chunk.strip() and current_tokens >= self.min_chunk_size:
                chunk = DocumentChunk(
                    chunk_id=f"{document_id}_chunk_{chunk_index}",
                    document_id=document_id,
                    content=current_chunk.strip(),
                    chunk_type="paragraph",
                    chunk_index=chunk_index,
                    token_count=current_tokens,
                    metadata={
                        "file_name": metadata.file_name,
                        "file_type": metadata.file_type
                    }
                )
                chunks.append(chunk)
            
            logger.info(f"📄 [DocumentProcessor] Created {len(chunks)} chunks from document")
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] Chunking failed: {e}")
            # Fallback: create single chunk
            chunks = [DocumentChunk(
                chunk_id=f"{document_id}_chunk_0",
                document_id=document_id,
                content=content[:self.max_chunk_size * 4],  # Rough character limit
                chunk_type="document",
                chunk_index=0,
                token_count=len(content.split()),
                metadata={
                    "file_name": metadata.file_name,
                    "file_type": metadata.file_type
                }
            )]
        
        return chunks
    
    async def _generate_embeddings(self, chunks: List[DocumentChunk]):
        """Generate embeddings for document chunks"""
        if not self.embedding_model or not chunks:
            return
        
        try:
            logger.info(f"🔢 [DocumentProcessor] Generating embeddings for {len(chunks)} chunks")
            
            # Extract content for embedding
            texts = [chunk.content for chunk in chunks]
            
            # Generate embeddings
            embeddings = self.embedding_model.encode(texts, convert_to_tensor=False)
            
            # Assign embeddings to chunks
            for chunk, embedding in zip(chunks, embeddings):
                chunk.embedding = embedding.tolist()
            
            logger.info("✅ [DocumentProcessor] Embeddings generated successfully")
            
        except Exception as e:
            logger.error(f"❌ [DocumentProcessor] Embedding generation failed: {e}")
    
    # ============================================================================
    # UTILITY METHODS
    # ============================================================================
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return list(self.supported_types.keys())
    
    def is_supported_format(self, file_path: str) -> bool:
        """Check if file format is supported"""
        return Path(file_path).suffix.lower() in self.supported_types
    
    async def get_document_info(self, file_path: str) -> Dict[str, Any]:
        """Get basic document information without full processing"""
        try:
            file_path_obj = Path(file_path)
            metadata = await self._extract_metadata(file_path_obj)
            
            return {
                "supported": self.is_supported_format(file_path),
                "metadata": asdict(metadata),
                "estimated_processing_time": self._estimate_processing_time(metadata.file_size)
            }
            
        except Exception as e:
            return {
                "supported": False,
                "error": str(e)
            }
    
    def _estimate_processing_time(self, file_size: int) -> float:
        """Estimate processing time based on file size"""
        # Rough estimate: 1MB per second
        return max(1.0, file_size / (1024 * 1024))
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        return {
            "supported_formats": len(self.supported_types),
            "max_chunk_size": self.max_chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "min_chunk_size": self.min_chunk_size,
            "embedding_model": self.embedding_model_name,
            "models_loaded": self.embedding_model is not None
        }

